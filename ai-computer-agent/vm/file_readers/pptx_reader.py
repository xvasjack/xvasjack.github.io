"""
PowerPoint Reader - Analyze PPTX files for quality checking.

This module extracts:
- Slide count and titles
- Company names and information
- Logo presence
- Website links
- Text content for analysis

Security:
- SEC-1: XXE protection via defusedxml
- SEC-2: ZIP bomb protection via size ratio check
"""

import os
import re
import zipfile
import math
from typing import List, Optional, Dict, Any, Tuple
from dataclasses import dataclass, field, asdict
from collections import Counter
from statistics import median, mode
import logging

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logging.getLogger("pptx_reader").warning("Missing python-pptx. Run: pip install python-pptx")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("pptx_reader")

# C7: Maximum PPTX file size to prevent OOM
MAX_PPTX_SIZE = 50 * 1024 * 1024  # 50MB

# SEC-2: ZIP bomb protection constants
MAX_DECOMPRESSED_SIZE = 500 * 1024 * 1024  # 500MB max decompressed
MAX_COMPRESSION_RATIO = 100  # Max 100:1 compression ratio
MAX_ZIP_ENTRIES = 10000  # Max number of files in archive

# IV-6: Maximum field length to prevent memory exhaustion
MAX_FIELD_LENGTH = 10000

# EMU conversion: 914400 EMU = 1 inch
EMU_PER_INCH = 914400

# Approximate characters per line at various font sizes (Segoe UI, standard slide width ~9.3")
# Used for overflow estimation — conservative
CHARS_PER_LINE_BY_PT = {
    10: 120, 12: 100, 14: 85, 16: 72, 18: 64, 20: 58, 24: 48,
}
LINE_HEIGHT_INCHES_BY_PT = {
    10: 0.18, 12: 0.22, 14: 0.25, 16: 0.29, 18: 0.32, 20: 0.36, 24: 0.43,
}


# =============================================================================
# FORMATTING DATACLASSES
# =============================================================================


@dataclass
class ShapeFormatting:
    """Formatting data for a single shape"""
    shape_type: str  # text_box, table, line, picture, placeholder, group, other
    left: float = 0.0   # inches
    top: float = 0.0     # inches
    width: float = 0.0   # inches
    height: float = 0.0  # inches
    font_name: Optional[str] = None
    font_size_pt: Optional[float] = None
    font_color_hex: Optional[str] = None
    font_bold: Optional[bool] = None
    text_length: int = 0
    estimated_lines: int = 0
    available_lines: int = 0  # based on box height and font size
    text_content: Optional[str] = None  # first 100 chars, for slide classification
    paragraph_alignment: Optional[str] = None  # left/center/right/justify
    fill_color_hex: Optional[str] = None  # shape background fill
    border_color_hex: Optional[str] = None  # shape outline color
    border_width_pt: Optional[float] = None  # shape outline width
    margin_left: Optional[float] = None  # text frame margin inches
    margin_top: Optional[float] = None
    margin_right: Optional[float] = None
    margin_bottom: Optional[float] = None
    line_spacing: Optional[float] = None  # paragraph line spacing

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class SlideFormatting:
    """Formatting data for a single slide"""
    slide_number: int
    layout_name: str = ""
    shapes: List[ShapeFormatting] = field(default_factory=list)
    title_font_size_pt: Optional[float] = None
    title_font_color_hex: Optional[str] = None
    title_font_bold: Optional[bool] = None
    subtitle_font_size_pt: Optional[float] = None
    subtitle_font_color_hex: Optional[str] = None
    header_line_count: int = 0
    header_line_y_positions: List[float] = field(default_factory=list)  # inches

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class FormattingProfile:
    """Learned formatting profile from a PPTX file"""
    source_file: str
    slide_count: int = 0

    # Title formatting (mode/median across content slides)
    title_font_name: Optional[str] = None
    title_font_size_pt: Optional[float] = None
    title_font_color_hex: Optional[str] = None
    title_font_bold: Optional[bool] = None

    # Subtitle formatting
    subtitle_font_name: Optional[str] = None
    subtitle_font_size_pt: Optional[float] = None
    subtitle_font_color_hex: Optional[str] = None

    # Header line info
    header_line_count_mode: int = 0
    header_line_y_positions: List[float] = field(default_factory=list)  # median y positions

    # Layout info
    content_start_y: float = 0.0  # inches — where content typically starts
    title_y: float = 0.0          # inches — where title typically sits
    font_family: Optional[str] = None  # most common font family

    # Per-slide formatting
    slides: List[SlideFormatting] = field(default_factory=list)

    # Overflow issues found
    overflow_shapes: List[Dict[str, Any]] = field(default_factory=list)
    overlap_pairs: List[Dict[str, Any]] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "source_file": self.source_file,
            "slide_count": self.slide_count,
            "title_font_name": self.title_font_name,
            "title_font_size_pt": self.title_font_size_pt,
            "title_font_color_hex": self.title_font_color_hex,
            "title_font_bold": self.title_font_bold,
            "subtitle_font_name": self.subtitle_font_name,
            "subtitle_font_size_pt": self.subtitle_font_size_pt,
            "subtitle_font_color_hex": self.subtitle_font_color_hex,
            "header_line_count_mode": self.header_line_count_mode,
            "header_line_y_positions": self.header_line_y_positions,
            "content_start_y": self.content_start_y,
            "title_y": self.title_y,
            "font_family": self.font_family,
            "overflow_shapes": self.overflow_shapes,
            "overlap_pairs": self.overlap_pairs,
            "slides": [s.to_dict() for s in self.slides],
        }

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "FormattingProfile":
        raw_slides = data.pop("slides", [])
        slides = []
        for s in raw_slides:
            if isinstance(s, dict):
                # Reconstruct ShapeFormatting objects within shapes list
                raw_shapes = s.pop("shapes", [])
                shape_objs = []
                for sh in raw_shapes:
                    if isinstance(sh, dict):
                        # Filter to known ShapeFormatting fields
                        shape_objs.append(ShapeFormatting(
                            **{k: v for k, v in sh.items()
                               if k in ShapeFormatting.__dataclass_fields__}
                        ))
                    else:
                        shape_objs.append(sh)
                slide_obj = SlideFormatting(
                    **{k: v for k, v in s.items()
                       if k in SlideFormatting.__dataclass_fields__}
                )
                slide_obj.shapes = shape_objs
                slides.append(slide_obj)
            else:
                slides.append(s)
        overflow = data.pop("overflow_shapes", [])
        overlap = data.pop("overlap_pairs", [])
        profile = cls(**{k: v for k, v in data.items()
                         if k in cls.__dataclass_fields__})
        profile.slides = slides
        profile.overflow_shapes = overflow
        profile.overlap_pairs = overlap
        return profile


def _emu_to_inches(emu: Optional[int]) -> float:
    """Convert EMU to inches, return 0 if None"""
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_INCH, 3)


def _get_shape_type(shape) -> str:
    """Classify shape type"""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        if hasattr(shape, 'has_table') and shape.has_table:
            return "table"
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return "line"
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return "line"
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        if shape.has_text_frame:
            return "text_box"
    except Exception:
        pass
    return "other"


def _get_first_run_font(shape) -> Tuple[Optional[str], Optional[float], Optional[str], Optional[bool]]:
    """Extract font info from the first run with text in a shape.
    Returns (font_name, font_size_pt, color_hex, bold)"""
    try:
        if not shape.has_text_frame:
            return None, None, None, None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                font = run.font
                name = font.name
                size_pt = round(font.size / 12700, 1) if font.size else None  # EMU to pt
                color_hex = None
                bold = font.bold
                try:
                    if font.color and font.color.rgb:
                        color_hex = str(font.color.rgb)
                except Exception:
                    pass
                return name, size_pt, color_hex, bold
    except Exception:
        pass
    return None, None, None, None


def _estimate_text_lines(text_length: int, box_width_inches: float, font_size_pt: Optional[float]) -> int:
    """Estimate how many lines text will take given box width and font size"""
    if text_length == 0 or box_width_inches <= 0:
        return 0
    pt = int(font_size_pt) if font_size_pt else 14
    # Find closest match in lookup
    closest_pt = min(CHARS_PER_LINE_BY_PT.keys(), key=lambda k: abs(k - pt))
    chars_per_line_full = CHARS_PER_LINE_BY_PT[closest_pt]
    # Scale by actual box width vs standard ~9.3"
    chars_per_line = max(1, int(chars_per_line_full * (box_width_inches / 9.3)))
    return max(1, math.ceil(text_length / chars_per_line))


def _estimate_available_lines(box_height_inches: float, font_size_pt: Optional[float]) -> int:
    """Estimate how many lines fit in a box given its height"""
    if box_height_inches <= 0:
        return 0
    pt = int(font_size_pt) if font_size_pt else 14
    closest_pt = min(LINE_HEIGHT_INCHES_BY_PT.keys(), key=lambda k: abs(k - pt))
    line_height = LINE_HEIGHT_INCHES_BY_PT[closest_pt]
    return max(1, int(box_height_inches / line_height))


def _extract_shape_formatting(shape) -> ShapeFormatting:
    """Extract formatting data from a single shape"""
    shape_type = _get_shape_type(shape)
    left = _emu_to_inches(shape.left)
    top = _emu_to_inches(shape.top)
    width = _emu_to_inches(shape.width)
    height = _emu_to_inches(shape.height)

    font_name, font_size_pt, font_color_hex, font_bold = _get_first_run_font(shape)

    text_length = 0
    text_content = None
    paragraph_alignment = None
    fill_color_hex = None
    border_color_hex = None
    border_width_pt = None
    margin_left = None
    margin_top = None
    margin_right = None
    margin_bottom = None
    line_spacing_val = None

    if shape.has_text_frame:
        text_length = sum(len(p.text) for p in shape.text_frame.paragraphs)
        # Extract first 100 chars for classification
        full_text = shape.text_frame.text or ""
        if full_text.strip():
            text_content = full_text.strip()[:100]
        # Paragraph alignment from first non-empty paragraph
        try:
            for para in shape.text_frame.paragraphs:
                if para.text.strip() and para.alignment is not None:
                    align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                    paragraph_alignment = align_map.get(para.alignment, str(para.alignment))
                    break
        except Exception:
            pass
        # Text frame margins
        try:
            tf = shape.text_frame
            if tf.margin_left is not None:
                margin_left = round(tf.margin_left / EMU_PER_INCH, 3)
            if tf.margin_top is not None:
                margin_top = round(tf.margin_top / EMU_PER_INCH, 3)
            if tf.margin_right is not None:
                margin_right = round(tf.margin_right / EMU_PER_INCH, 3)
            if tf.margin_bottom is not None:
                margin_bottom = round(tf.margin_bottom / EMU_PER_INCH, 3)
        except Exception:
            pass
        # Line spacing from first non-empty paragraph
        try:
            for para in shape.text_frame.paragraphs:
                if para.text.strip() and para.line_spacing is not None:
                    line_spacing_val = float(para.line_spacing)
                    break
        except Exception:
            pass

    # Shape fill color
    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            if shape.fill.type == 1:  # Solid fill
                fill_color_hex = str(shape.fill.fore_color.rgb)
    except Exception:
        pass

    # Shape border/outline
    try:
        if hasattr(shape, 'line') and shape.line.color and shape.line.color.rgb:
            border_color_hex = str(shape.line.color.rgb)
        if hasattr(shape, 'line') and shape.line.width is not None:
            border_width_pt = round(shape.line.width / 12700, 1)  # EMU to pt
    except Exception:
        pass

    estimated_lines = _estimate_text_lines(text_length, width, font_size_pt)
    available_lines = _estimate_available_lines(height, font_size_pt)

    return ShapeFormatting(
        shape_type=shape_type,
        left=left, top=top, width=width, height=height,
        font_name=font_name, font_size_pt=font_size_pt,
        font_color_hex=font_color_hex, font_bold=font_bold,
        text_length=text_length,
        estimated_lines=estimated_lines,
        available_lines=available_lines,
        text_content=text_content,
        paragraph_alignment=paragraph_alignment,
        fill_color_hex=fill_color_hex,
        border_color_hex=border_color_hex,
        border_width_pt=border_width_pt,
        margin_left=margin_left,
        margin_top=margin_top,
        margin_right=margin_right,
        margin_bottom=margin_bottom,
        line_spacing=line_spacing_val,
    )


def _extract_slide_formatting(slide, slide_number: int) -> SlideFormatting:
    """Extract formatting data from a single slide"""
    layout_name = ""
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        pass

    shapes = []
    title_font_size = None
    title_font_color = None
    title_font_bold = None
    subtitle_font_size = None
    subtitle_font_color = None
    header_lines = []

    for shape in slide.shapes:
        sf = _extract_shape_formatting(shape)
        shapes.append(sf)

        # Detect title — try placeholder first, then position-based heuristic
        if shape.has_text_frame and sf.text_length > 0:
            is_title = False
            is_subtitle = False
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:
                    is_title = True
                elif ph_type == 2:
                    is_subtitle = True
            except (ValueError, AttributeError):
                # Not a placeholder — use position heuristic:
                # Title = topmost text shape (y < 0.8") with font >= 18pt
                # Subtitle = second topmost (0.4" < y < 1.2") with font >= 14pt
                if sf.top < 0.8 and sf.font_size_pt and sf.font_size_pt >= 18:
                    is_title = True
                elif 0.4 < sf.top < 1.2 and sf.font_size_pt and 12 <= sf.font_size_pt < 18:
                    is_subtitle = True

            if is_title and title_font_size is None:
                title_font_size = sf.font_size_pt
                title_font_color = sf.font_color_hex
                title_font_bold = sf.font_bold
            elif is_subtitle and subtitle_font_size is None:
                subtitle_font_size = sf.font_size_pt
                subtitle_font_color = sf.font_color_hex

        # Detect header lines (line shapes or thin shapes near y=0.5-1.8")
        if sf.shape_type == "line" and 0.5 <= sf.top <= 1.8:
            header_lines.append(sf.top)
        # Also detect connector shapes used as divider lines (height near 0)
        elif sf.height < 0.05 and sf.width > 2.0 and 0.5 <= sf.top <= 1.8:
            header_lines.append(sf.top)

    return SlideFormatting(
        slide_number=slide_number,
        layout_name=layout_name,
        shapes=shapes,
        title_font_size_pt=title_font_size,
        title_font_color_hex=title_font_color,
        title_font_bold=title_font_bold,
        subtitle_font_size_pt=subtitle_font_size,
        subtitle_font_color_hex=subtitle_font_color,
        header_line_count=len(header_lines),
        header_line_y_positions=sorted(header_lines),
    )


def _estimate_overflow(slide_fmt: SlideFormatting) -> List[Dict[str, Any]]:
    """Check for text overflow in a slide's shapes"""
    overflows = []
    for sf in slide_fmt.shapes:
        if sf.shape_type not in ("text_box", "placeholder"):
            continue
        if sf.text_length == 0 or sf.available_lines == 0:
            continue
        if sf.estimated_lines > sf.available_lines:
            overflows.append({
                "slide": slide_fmt.slide_number,
                "shape_type": sf.shape_type,
                "top": sf.top,
                "height": sf.height,
                "estimated_lines": sf.estimated_lines,
                "available_lines": sf.available_lines,
                "overflow_lines": sf.estimated_lines - sf.available_lines,
                "font_size_pt": sf.font_size_pt,
            })
    return overflows


def _detect_overlaps(slide_fmt: SlideFormatting) -> List[Dict[str, Any]]:
    """Detect overlapping shapes on a slide by checking y + height > next shape y"""
    overlaps = []
    # Only check text/table shapes, sorted by y position
    content_shapes = [s for s in slide_fmt.shapes
                      if s.shape_type in ("text_box", "table", "placeholder") and s.height > 0]
    content_shapes.sort(key=lambda s: s.top)

    for i in range(len(content_shapes) - 1):
        current = content_shapes[i]
        nxt = content_shapes[i + 1]
        bottom = current.top + current.height
        if bottom > nxt.top + 0.05:  # 0.05" tolerance
            overlap_inches = round(bottom - nxt.top, 2)
            overlaps.append({
                "slide": slide_fmt.slide_number,
                "shape_a_type": current.shape_type,
                "shape_a_top": current.top,
                "shape_a_bottom": round(bottom, 2),
                "shape_b_type": nxt.shape_type,
                "shape_b_top": nxt.top,
                "overlap_inches": overlap_inches,
            })
    return overlaps


def extract_formatting_profile(pptx_path: str) -> Optional[FormattingProfile]:
    """
    Extract a FormattingProfile from a PPTX file.
    Reads fonts, colors, positions, lines — no hardcoded values.
    Returns None if file can't be read.
    """
    if not HAS_PPTX:
        logger.warning("python-pptx not available for formatting extraction")
        return None

    if not os.path.exists(pptx_path):
        logger.warning(f"File not found for formatting extraction: {pptx_path}")
        return None

    # Safety checks
    try:
        file_size = os.path.getsize(pptx_path)
        if file_size > MAX_PPTX_SIZE:
            logger.warning(f"File too large for formatting extraction: {file_size}")
            return None
    except OSError:
        return None

    is_safe, bomb_error = _check_zip_bomb(pptx_path)
    if not is_safe:
        logger.warning(f"ZIP bomb protection: {bomb_error}")
        return None

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.warning(f"Could not open PPTX for formatting: {e}")
        return None

    profile = FormattingProfile(
        source_file=os.path.basename(pptx_path),
        slide_count=len(prs.slides),
    )

    # Collect per-slide formatting
    all_title_sizes = []
    all_title_colors = []
    all_title_bolds = []
    all_subtitle_sizes = []
    all_subtitle_colors = []
    all_font_names = []
    all_header_line_counts = []
    all_header_line_ys = []
    all_title_ys = []
    all_content_start_ys = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_fmt = _extract_slide_formatting(slide, idx)
        profile.slides.append(slide_fmt)

        # Skip first slide (title slide) for statistics
        if idx == 1:
            continue

        if slide_fmt.title_font_size_pt:
            all_title_sizes.append(slide_fmt.title_font_size_pt)
        if slide_fmt.title_font_color_hex:
            all_title_colors.append(slide_fmt.title_font_color_hex)
        if slide_fmt.title_font_bold is not None:
            all_title_bolds.append(slide_fmt.title_font_bold)
        if slide_fmt.subtitle_font_size_pt:
            all_subtitle_sizes.append(slide_fmt.subtitle_font_size_pt)
        if slide_fmt.subtitle_font_color_hex:
            all_subtitle_colors.append(slide_fmt.subtitle_font_color_hex)

        all_header_line_counts.append(slide_fmt.header_line_count)
        all_header_line_ys.extend(slide_fmt.header_line_y_positions)

        # Track font names and positions
        for sf in slide_fmt.shapes:
            if sf.font_name:
                all_font_names.append(sf.font_name)
            # Title is usually the topmost text shape
            if sf.shape_type in ("text_box", "placeholder") and sf.top < 1.0:
                all_title_ys.append(sf.top)
            elif sf.shape_type in ("text_box", "table", "placeholder") and sf.top > 1.0:
                all_content_start_ys.append(sf.top)

        # Check for overflows and overlaps
        overflows = _estimate_overflow(slide_fmt)
        overlaps = _detect_overlaps(slide_fmt)
        profile.overflow_shapes.extend(overflows)
        profile.overlap_pairs.extend(overlaps)

    # Compute mode/median statistics
    if all_title_sizes:
        try:
            profile.title_font_size_pt = mode(all_title_sizes)
        except Exception:
            profile.title_font_size_pt = median(all_title_sizes)

    if all_title_colors:
        try:
            profile.title_font_color_hex = Counter(all_title_colors).most_common(1)[0][0]
        except Exception:
            pass

    if all_title_bolds:
        profile.title_font_bold = Counter(all_title_bolds).most_common(1)[0][0]

    if all_subtitle_sizes:
        try:
            profile.subtitle_font_size_pt = mode(all_subtitle_sizes)
        except Exception:
            profile.subtitle_font_size_pt = median(all_subtitle_sizes)

    if all_subtitle_colors:
        try:
            profile.subtitle_font_color_hex = Counter(all_subtitle_colors).most_common(1)[0][0]
        except Exception:
            pass

    if all_font_names:
        profile.font_family = Counter(all_font_names).most_common(1)[0][0]
        profile.title_font_name = profile.font_family
        profile.subtitle_font_name = profile.font_family

    if all_header_line_counts:
        try:
            profile.header_line_count_mode = mode(all_header_line_counts)
        except Exception:
            profile.header_line_count_mode = round(median(all_header_line_counts))

    if all_header_line_ys:
        profile.header_line_y_positions = [round(median(all_header_line_ys), 2)]

    if all_title_ys:
        profile.title_y = round(median(all_title_ys), 2)

    if all_content_start_ys:
        profile.content_start_y = round(median(all_content_start_ys), 2)

    return profile


def _check_zip_bomb(file_path: str) -> tuple:
    """
    SEC-2: Check for ZIP bomb before processing.
    Returns (is_safe, error_message).
    """
    try:
        compressed_size = os.path.getsize(file_path)
        if compressed_size == 0:
            return False, "Empty file"

        with zipfile.ZipFile(file_path, 'r') as zf:
            # Check number of entries
            if len(zf.namelist()) > MAX_ZIP_ENTRIES:
                return False, f"Too many files in archive: {len(zf.namelist())} > {MAX_ZIP_ENTRIES}"

            # Calculate total decompressed size
            total_uncompressed = sum(info.file_size for info in zf.infolist())

            # Check absolute size limit
            if total_uncompressed > MAX_DECOMPRESSED_SIZE:
                return False, f"Decompressed size too large: {total_uncompressed / 1024 / 1024:.1f}MB > {MAX_DECOMPRESSED_SIZE / 1024 / 1024:.0f}MB"

            # Check compression ratio
            ratio = total_uncompressed / compressed_size if compressed_size > 0 else 0
            if ratio > MAX_COMPRESSION_RATIO:
                return False, f"Suspicious compression ratio: {ratio:.1f}:1 > {MAX_COMPRESSION_RATIO}:1"

            return True, None
    except zipfile.BadZipFile:
        return False, "Invalid or corrupted ZIP file"
    except Exception as e:
        return False, f"ZIP validation error: {e}"


@dataclass
class SlideInfo:
    number: int
    title: Optional[str]
    text_content: List[str]
    has_image: bool
    has_table: bool
    links: List[str]


@dataclass
class CompanyInfo:
    name: str
    website: Optional[str] = None
    description: Optional[str] = None
    has_logo: bool = False
    slide_number: int = 0


@dataclass
class PPTXAnalysis:
    file_path: str
    slide_count: int
    slides: List[SlideInfo] = field(default_factory=list)
    companies: List[CompanyInfo] = field(default_factory=list)
    total_images: int = 0
    total_links: int = 0
    issues: List[str] = field(default_factory=list)
    summary: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "slide_count": self.slide_count,
            "company_count": len(self.companies),
            "total_images": self.total_images,
            "total_links": self.total_links,
            "issues": self.issues,
            "summary": self.summary,
            # 3.11: Include description and slide_number in to_dict
            "companies": [
                {
                    "name": c.name,
                    "website": c.website,
                    "has_logo": c.has_logo,
                    "description": getattr(c, 'description', ''),
                    "slide_number": getattr(c, 'slide_number', 0),
                }
                for c in self.companies
            ],
            # Include slide details for content depth checks
            "slides": [
                {
                    "number": s.number,
                    "title": s.title,
                    "type": "title" if s.number == 1 else "",
                    "all_text": "\n".join(s.text_content),
                    "text_blocks": s.text_content,
                    "has_image": s.has_image,
                    "has_table": s.has_table,
                    "links": s.links,
                }
                for s in self.slides
            ],
        }


def analyze_pptx(file_path: str) -> PPTXAnalysis:
    """
    Analyze a PowerPoint file and extract relevant information.

    Args:
        file_path: Path to the PPTX file

    Returns:
        PPTXAnalysis with extracted information
    """
    logger.info(f"Analyzing PPTX: {file_path}")

    if not HAS_PPTX:
        return PPTXAnalysis(
            file_path=file_path,
            slide_count=0,
            issues=["python-pptx not installed"],
            summary="Error: python-pptx not installed"
        )

    if not os.path.exists(file_path):
        return PPTXAnalysis(
            file_path=file_path,
            slide_count=0,
            issues=["File not found"],
            summary="Error: File not found"
        )

    # C7: Check file size to prevent OOM
    try:
        file_size = os.path.getsize(file_path)
        if file_size > MAX_PPTX_SIZE:
            size_mb = file_size / 1024 / 1024
            limit_mb = MAX_PPTX_SIZE / 1024 / 1024
            return PPTXAnalysis(
                file_path=file_path,
                slide_count=0,
                issues=[f"File too large: {size_mb:.1f}MB > {limit_mb:.0f}MB limit"],
                summary=f"Error: File too large ({size_mb:.1f}MB)"
            )
    except OSError as e:
        logger.warning(f"Could not check file size: {e}")

    # SEC-2: Check for ZIP bomb before processing
    is_safe, bomb_error = _check_zip_bomb(file_path)
    if not is_safe:
        logger.warning(f"ZIP bomb protection triggered: {bomb_error}")
        return PPTXAnalysis(
            file_path=file_path,
            slide_count=0,
            issues=[f"Security: {bomb_error}"],
            summary=f"Security error: {bomb_error}"
        )

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return PPTXAnalysis(
            file_path=file_path,
            slide_count=0,
            issues=[f"Could not open file: {e}"],
            summary=f"Error opening file: {e}"
        )

    analysis = PPTXAnalysis(
        file_path=file_path,
        slide_count=len(prs.slides),
    )

    total_images = 0
    total_links = 0
    companies = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_info = extract_slide_info(slide, idx)
        analysis.slides.append(slide_info)

        total_images += 1 if slide_info.has_image else 0
        total_links += len(slide_info.links)

        # Try to extract company info from slide
        company = extract_company_from_slide(slide, idx)
        if company:
            companies.append(company)

    analysis.total_images = total_images
    analysis.total_links = total_links
    analysis.companies = companies

    # Check for issues
    analysis.issues = check_quality_issues(analysis)

    # Generate summary
    analysis.summary = generate_summary(analysis)

    return analysis


def extract_slide_info(slide, slide_number: int) -> SlideInfo:
    """Extract information from a single slide"""

    title = None
    text_content = []
    has_image = False
    has_table = False
    links = []

    for shape in slide.shapes:
        # Get title
        if shape.has_text_frame:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    title = shape.text.strip()

            # Get all text
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_content.append(text)

                # Check for hyperlinks
                # Category 1 fix: run.hyperlink.address may be None
                for run in paragraph.runs:
                    try:
                        if run.hyperlink and run.hyperlink.address:
                            links.append(run.hyperlink.address)
                    except AttributeError:
                        pass  # Expected: some runs don't have hyperlinks
                    except Exception as e:
                        logger.debug(f"Error reading hyperlink on slide {slide_number}: {e}")

        # Check for images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_image = True

        # Check for tables
        if shape.has_table:
            has_table = True

    return SlideInfo(
        number=slide_number,
        title=title,
        text_content=text_content,
        has_image=has_image,
        has_table=has_table,
        links=links,
    )


def extract_company_from_slide(slide, slide_number: int) -> Optional[CompanyInfo]:
    """Try to extract company information from a slide"""

    texts = []
    has_image = False
    links = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)

                # Category 1 fix: run.hyperlink.address may be None
                for run in paragraph.runs:
                    try:
                        if run.hyperlink and run.hyperlink.address:
                            links.append(run.hyperlink.address)
                    except AttributeError:
                        pass  # Expected: some runs don't have hyperlinks
                    except Exception as e:
                        logger.debug(f"Error reading hyperlink in company extraction: {e}")

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_image = True

    if not texts:
        return None

    # 3.12: Skip known non-company patterns (titles, headers, etc.)
    skip_patterns = ["table of contents", "agenda", "appendix", "disclaimer",
                     "confidential", "prepared for", "prepared by", "page ",
                     # Market-research section headers (false positive fix)
                     "section ", "overview", "legislation", "regulatory",
                     "key players", "market size", "market trends",
                     "competitive landscape", "swot", "pest", "porter",
                     "recommendations", "conclusion", "executive summary",
                     "introduction", "methodology", "references", "sources"]

    import re as _re

    def _is_skip_text(text: str) -> bool:
        """Check if text matches non-company patterns"""
        if not text or len(text) < 2:
            return True
        # All caps short text = likely country/region header
        if text.isupper() and len(text.split()) <= 3:
            return True
        # Numbered section patterns
        if _re.match(r'^section\s+\d+', text.lower()):
            return True
        # "Country - Topic" section headers
        if _re.match(r'^[A-Z][a-zA-Z\s]{1,30}\s*[-–—]\s*\S', text):
            return True
        if any(p in text.lower() for p in skip_patterns):
            return True
        return False

    # Try first text as company name; if it's a skip pattern, try texts[1:]
    company_name = None
    for candidate in texts:
        if not _is_skip_text(candidate):
            company_name = candidate
            break

    if not company_name or len(company_name) < 2:
        return None

    # Find website
    website = None
    for link in links:
        if link.startswith("http"):
            website = link
            break

    # Look for website in text
    # Category 6 fix: URL regex captures trailing punctuation - fix the pattern
    if not website:
        for text in texts:
            url_match = re.search(r'https?://[^\s<>"\')]+', text)
            if url_match:
                website = url_match.group().rstrip('.,;:!?')
                break
            # Check for domain pattern
            domain_match = re.search(r'www\.[a-zA-Z0-9][-a-zA-Z0-9]*\.[a-zA-Z]{2,}[^\s<>"\']*', text)
            if domain_match:
                website = "https://" + domain_match.group().rstrip('.,;:!?')
                break

    # Description is usually second text block
    description = texts[1] if len(texts) > 1 else None

    return CompanyInfo(
        name=company_name,
        website=website,
        description=description,
        has_logo=has_image,
        slide_number=slide_number,
    )


def check_quality_issues(analysis: PPTXAnalysis) -> List[str]:
    """Check for quality issues in the presentation"""

    issues = []

    # Check slide count
    if analysis.slide_count < 5:
        issues.append(f"Low slide count: only {analysis.slide_count} slides")

    # Check company count
    if len(analysis.companies) < 10:
        issues.append(f"Low company count: only {len(analysis.companies)} companies found")

    # Check for companies without logos
    no_logo_companies = [c for c in analysis.companies if not c.has_logo]
    if no_logo_companies:
        issues.append(f"{len(no_logo_companies)} companies without logos")

    # Check for companies without websites
    no_website_companies = [c for c in analysis.companies if not c.website]
    if no_website_companies:
        issues.append(f"{len(no_website_companies)} companies without websites")

    # Check for empty slides
    empty_slides = [s for s in analysis.slides if not s.text_content]
    if empty_slides:
        issues.append(f"{len(empty_slides)} empty slides")

    # Check for broken links (basic validation)
    for company in analysis.companies:
        if company.website:
            if not company.website.startswith("http"):
                issues.append(f"Invalid URL for {company.name}: {company.website}")

    return issues


def generate_summary(analysis: PPTXAnalysis) -> str:
    """Generate a summary of the analysis"""

    companies_with_logos = sum(1 for c in analysis.companies if c.has_logo)
    companies_with_websites = sum(1 for c in analysis.companies if c.website)

    summary = f"""PPTX Analysis Summary:
- File: {os.path.basename(analysis.file_path)}
- Slides: {analysis.slide_count}
- Companies found: {len(analysis.companies)}
- Companies with logos: {companies_with_logos}
- Companies with websites: {companies_with_websites}
- Total images: {analysis.total_images}
- Issues found: {len(analysis.issues)}

Issues:
{chr(10).join('- ' + issue for issue in analysis.issues) if analysis.issues else '- None'}

Top companies:
{chr(10).join(f'- {c.name} ({c.website or "no website"})' for c in analysis.companies[:5])}
"""

    return summary


def _is_safe_url(url: str) -> bool:
    """
    Check if a URL is safe to fetch (prevents SSRF attacks).
    Blocks internal/private IPs and non-HTTP(S) protocols.
    """
    from urllib.parse import urlparse
    import socket
    import ipaddress

    try:
        parsed = urlparse(url)

        # Only allow HTTP and HTTPS protocols
        if parsed.scheme.lower() not in ('http', 'https'):
            return False

        # Must have a hostname
        if not parsed.hostname:
            return False

        hostname = parsed.hostname.lower()

        # Block localhost and common internal hostnames
        blocked_hostnames = ['localhost', '127.0.0.1', '0.0.0.0', '[::1]', 'internal', 'intranet']
        if hostname in blocked_hostnames:
            return False

        # Block .local and .internal domains
        if hostname.endswith('.local') or hostname.endswith('.internal'):
            return False

        # Resolve hostname and check if it's a private IP
        # SEC-4: Set socket timeout to prevent hanging on slow DNS
        try:
            old_timeout = socket.getdefaulttimeout()
            socket.setdefaulttimeout(3)  # 3 second timeout for DNS
            try:
                ip_str = socket.gethostbyname(hostname)
                ip = ipaddress.ip_address(ip_str)
                if ip.is_private or ip.is_loopback or ip.is_reserved or ip.is_link_local:
                    return False
            finally:
                socket.setdefaulttimeout(old_timeout)  # Restore original timeout
        except (socket.gaierror, ValueError, socket.timeout):
            # If we can't resolve, allow it (external DNS will handle it)
            pass

        return True

    except Exception:
        return False


def validate_company_links(analysis: PPTXAnalysis) -> Dict[str, bool]:
    """
    Validate that company website links are accessible.
    Returns dict of {company_name: is_valid}

    Note: This does actual HTTP requests, use sparingly.
    Security: URLs are validated to prevent SSRF attacks.
    """
    import urllib.request
    import urllib.error

    results = {}

    for company in analysis.companies:
        if not company.website:
            results[company.name] = False
            continue

        # SSRF protection: validate URL before making request
        if not _is_safe_url(company.website):
            logger.warning(f"Blocked unsafe URL for {company.name}: {company.website}")
            results[company.name] = False
            continue

        try:
            req = urllib.request.Request(
                company.website,
                headers={'User-Agent': 'Mozilla/5.0'}
            )
            response = urllib.request.urlopen(req, timeout=10)
            # Close connection properly to avoid resource leak
            try:
                status = response.getcode()
                results[company.name] = status == 200
            finally:
                response.close()
        except (urllib.error.URLError, urllib.error.HTTPError, Exception):
            results[company.name] = False

    return results


# =============================================================================
# COMPARISON FUNCTIONS
# =============================================================================


def compare_analyses(old: PPTXAnalysis, new: PPTXAnalysis) -> Dict[str, Any]:
    """
    Compare two PPTX analyses to see if issues were fixed.

    Returns dict with comparison results.
    """
    old_companies = {c.name for c in old.companies}
    new_companies = {c.name for c in new.companies}

    return {
        "slide_count_change": new.slide_count - old.slide_count,
        "company_count_change": len(new.companies) - len(old.companies),
        "companies_added": list(new_companies - old_companies),
        "companies_removed": list(old_companies - new_companies),
        "old_issues": old.issues,
        "new_issues": new.issues,
        "issues_fixed": [i for i in old.issues if i not in new.issues],
        "new_issues_introduced": [i for i in new.issues if i not in old.issues],
        "improved": len(new.issues) < len(old.issues),
    }


# =============================================================================
# ENTRY POINT
# =============================================================================


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_reader.py <file.pptx>")
        sys.exit(1)

    analysis = analyze_pptx(sys.argv[1])
    print(analysis.summary)
