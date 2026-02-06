#!/usr/bin/env python3
"""
Pattern Extractor for PowerPoint Templates

This script analyzes a reference PowerPoint template file and extracts 12 distinct
layout patterns based on shape positions, sizes, fonts, colors, and other properties.

The extracted patterns are saved to a JSON file that can be used by the slide
generator to create new presentations with consistent styling.
"""

import argparse
import json
import logging
import sys
from collections import Counter
from pathlib import Path
from typing import Dict, List, Any, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# EMU to inches conversion constant (1 inch = 914400 EMU)
EMU_PER_INCH = 914400

# Pattern classification mapping: slide numbers to pattern names
PATTERN_MAPPING = {
    1: "cover",
    2: "toc_divider",
    3: "matrix_2x2",           # executive summary
    4: "matrix_2x2",           # Vietnam opportunities overview
    5: "toc_divider",
    6: "label_row_table",      # regulatory summary
    7: "data_table_reference", # foundational acts
    8: "text_policy_block",    # policies
    9: "text_policy_block",    # regulations
    10: "text_policy_block",   # incentives
    11: "toc_divider",
    12: "label_row_table",     # market summary
    13: "chart_insight_panels",# TPES
    14: "chart_insight_panels",# energy consumption
    15: "chart_insight_panels",# electricity
    16: "chart_callout_boxes", # electricity price
    17: "chart_callout_boxes", # natural gas
    18: "chart_callout_boxes", # LNG infrastructure
    19: "chart_callout_boxes", # gas import
    20: "toc_divider",
    21: "label_row_table",     # Japanese players summary
    22: "data_table_highlighted", # player presence
    23: "case_study_rows",     # Osaka Gas case 1
    24: "case_study_rows",     # Osaka Gas case 2
    25: "case_study_rows",     # Japanese PR reference
    26: "dual_chart_financial",# SOGEC financials
    27: "case_study_rows",     # Toho Gas case 1
    28: "case_study_rows",     # Toho Gas case 2
    29: "dual_chart_financial",# PSM financials
    30: "toc_divider",
    31: "chart_callout_boxes", # electricity supply
    32: "chart_callout_boxes", # electricity price
    33: "glossary_table",
    34: "glossary_table",
}

# Pattern descriptions and use cases
PATTERN_INFO = {
    "cover": {
        "description": "Title slide with main heading, subtitle, and branding",
        "bestFor": ["presentation opening", "title page", "cover slide"]
    },
    "toc_divider": {
        "description": "Table of contents or section divider with navigation",
        "bestFor": ["agenda", "section breaks", "navigation slides"]
    },
    "matrix_2x2": {
        "description": "2x2 grid layout for quadrant analysis or comparison",
        "bestFor": ["SWOT analysis", "four-quadrant frameworks", "comparison matrices"]
    },
    "structured_table": {
        "description": "Label-row table with headers and structured data rows",
        "bestFor": ["feature lists", "specifications", "structured data"]
    },
    "multicolumn_table": {
        "description": "Multi-column data table with headers and detailed content",
        "bestFor": ["detailed comparisons", "data tables", "multi-attribute listings"]
    },
    "text_policy_block": {
        "description": "Text-heavy layout for policies, definitions, or detailed explanations",
        "bestFor": ["policy text", "detailed descriptions", "methodology explanations"]
    },
    "chart_insight_panels": {
        "description": "Chart with insight panels showing key findings and analysis",
        "bestFor": ["data visualization with insights", "analytical findings", "key metrics"]
    },
    "chart_callout_boxes": {
        "description": "Chart with callout boxes for annotations and highlights",
        "bestFor": ["annotated charts", "highlighted data points", "visual storytelling"]
    },
    "case_study_rows": {
        "description": "Row-based layout for case studies with images and descriptions",
        "bestFor": ["case studies", "examples", "portfolio items"]
    },
    "dual_chart_financial": {
        "description": "Side-by-side charts for financial or comparative data",
        "bestFor": ["financial comparisons", "before/after", "dual metrics"]
    },
    "diagram_text_split": {
        "description": "Split layout with diagram on one side and text on the other",
        "bestFor": ["process diagrams", "workflow explanations", "concept illustrations"]
    },
    "glossary_table": {
        "description": "Two-column table for terms and definitions",
        "bestFor": ["glossary", "definitions", "term explanations"]
    },
    "label_row_table": {
        "description": "Structured label-row table with themed category labels on left",
        "bestFor": ["section summary", "themed overview", "category breakdown"]
    },
    "data_table_reference": {
        "description": "Multi-column reference data table for structured comparative data",
        "bestFor": ["regulation list", "law reference", "structured data"]
    },
    "data_table_highlighted": {
        "description": "Multi-column data table with colored left border for highlighting",
        "bestFor": ["company comparison", "competitive analysis", "highlighted data"]
    }
}


def emu_to_inches(emu: int) -> float:
    """Convert EMU (English Metric Units) to inches.

    Args:
        emu: Value in EMU units

    Returns:
        Value in inches, rounded to 2 decimal places
    """
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_INCH, 2)


def extract_color(color_obj) -> Optional[str]:
    """Extract color from a pptx color object.

    Args:
        color_obj: pptx color object (FillFormat, ColorFormat, etc.)

    Returns:
        Hex color string (e.g., "#1F4E78") or None
    """
    try:
        if hasattr(color_obj, 'rgb'):
            rgb = color_obj.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
        elif hasattr(color_obj, 'color') and hasattr(color_obj.color, 'rgb'):
            rgb = color_obj.color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
    except Exception as e:
        logger.debug(f"Could not extract color: {e}")
    return None


def extract_font_props(font) -> Dict[str, Any]:
    """Extract font properties from a font object.

    Args:
        font: pptx Font object

    Returns:
        Dictionary with font properties
    """
    props = {}

    try:
        if font.name:
            props['family'] = font.name
    except Exception:
        pass

    try:
        if font.size:
            props['size'] = round(font.size.pt, 1)
    except Exception:
        pass

    try:
        if font.color:
            color = extract_color(font.color)
            if color:
                props['color'] = color
    except Exception:
        pass

    try:
        props['bold'] = bool(font.bold)
    except Exception:
        props['bold'] = False

    try:
        props['italic'] = bool(font.italic)
    except Exception:
        props['italic'] = False

    return props


def extract_shape_data(shape) -> Optional[Dict[str, Any]]:
    """Extract data from a single shape.

    Args:
        shape: pptx Shape object

    Returns:
        Dictionary with shape properties or None if shape should be skipped
    """
    try:
        shape_data = {
            'type': shape.shape_type.name if hasattr(shape, 'shape_type') else 'UNKNOWN',
            'x': emu_to_inches(shape.left),
            'y': emu_to_inches(shape.top),
            'width': emu_to_inches(shape.width),
            'height': emu_to_inches(shape.height),
        }

        # Add shape name if available
        if hasattr(shape, 'name') and shape.name:
            shape_data['name'] = shape.name

        # Extract text content and formatting
        if hasattr(shape, 'text') and shape.text:
            shape_data['text'] = shape.text.strip()

        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            if text_frame.paragraphs:
                fonts = []
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_props = extract_font_props(run.font)
                        if font_props:
                            fonts.append(font_props)

                if fonts:
                    shape_data['fonts'] = fonts

        # Extract fill color
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:  # Solid fill
                    color = extract_color(shape.fill.fore_color)
                    if color:
                        shape_data['fill_color'] = color
            except Exception:
                pass

        # Extract line/border properties
        if hasattr(shape, 'line'):
            try:
                line_data = {}
                if shape.line.color:
                    color = extract_color(shape.line.color)
                    if color:
                        line_data['color'] = color
                if shape.line.width:
                    line_data['width'] = emu_to_inches(shape.line.width)
                if line_data:
                    shape_data['line'] = line_data
            except Exception:
                pass

        # Extract table data
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_data = extract_table_data(shape.table)
            if table_data:
                shape_data['table'] = table_data

        # Extract chart data
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_data = extract_chart_data(shape.chart)
            if chart_data:
                shape_data['chart'] = chart_data

        return shape_data

    except Exception as e:
        logger.debug(f"Error extracting shape data: {e}")
        return None


def extract_table_data(table) -> Dict[str, Any]:
    """Extract data from a table shape.

    Args:
        table: pptx Table object

    Returns:
        Dictionary with table properties
    """
    try:
        table_data = {
            'rows': len(table.rows),
            'columns': len(table.columns),
            'column_widths': [emu_to_inches(col.width) for col in table.columns],
            'row_heights': [emu_to_inches(row.height) for row in table.rows],
        }

        # Extract header row formatting (first row)
        if table.rows:
            header_cells = []
            for cell in table.rows[0].cells:
                cell_data = {}
                if cell.text:
                    cell_data['text'] = cell.text.strip()

                # Extract cell fill color
                if hasattr(cell.fill, 'fore_color'):
                    color = extract_color(cell.fill.fore_color)
                    if color:
                        cell_data['fill_color'] = color

                # Extract text formatting
                if cell.text_frame and cell.text_frame.paragraphs:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font_props = extract_font_props(run.font)
                            if font_props:
                                cell_data['font'] = font_props
                                break
                        if 'font' in cell_data:
                            break

                header_cells.append(cell_data)

            if header_cells:
                table_data['header_cells'] = header_cells

        return table_data

    except Exception as e:
        logger.debug(f"Error extracting table data: {e}")
        return {}


def extract_chart_data(chart) -> Dict[str, Any]:
    """Extract data from a chart shape.

    Args:
        chart: pptx Chart object

    Returns:
        Dictionary with chart properties
    """
    try:
        chart_data = {
            'chart_type': chart.chart_type.name if hasattr(chart, 'chart_type') else 'UNKNOWN',
        }

        # Extract chart colors from series
        if hasattr(chart, 'plots') and chart.plots:
            colors = []
            for plot in chart.plots:
                if hasattr(plot, 'series'):
                    for series in plot.series:
                        try:
                            if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                                color = extract_color(series.format.fill.fore_color)
                                if color:
                                    colors.append(color)
                        except Exception:
                            pass

            if colors:
                chart_data['colors'] = colors

        return chart_data

    except Exception as e:
        logger.debug(f"Error extracting chart data: {e}")
        return {}


def extract_slide_shapes(slide) -> List[Dict[str, Any]]:
    """Extract all shapes from a slide.

    Args:
        slide: pptx Slide object

    Returns:
        List of shape dictionaries
    """
    shapes = []

    for shape in slide.shapes:
        shape_data = extract_shape_data(shape)
        if shape_data:
            shapes.append(shape_data)

    return shapes


def classify_slide(slide_num: int) -> Optional[str]:
    """Classify a slide into a pattern based on its slide number.

    Args:
        slide_num: Slide number (1-indexed)

    Returns:
        Pattern name or None if not in mapping
    """
    return PATTERN_MAPPING.get(slide_num)


def extract_global_style(all_slides_data: List[Dict[str, Any]], prs: Presentation) -> Dict[str, Any]:
    """Extract global style constants from all slides.

    Args:
        all_slides_data: List of slide data dictionaries
        prs: Presentation object

    Returns:
        Dictionary with global style properties
    """
    # Collect all fonts and colors
    all_fonts = []
    all_colors = []
    all_fill_colors = []

    for slide_data in all_slides_data:
        for shape in slide_data.get('shapes', []):
            # Collect fonts
            if 'fonts' in shape:
                all_fonts.extend(shape['fonts'])

            # Collect colors
            if 'fill_color' in shape:
                all_fill_colors.append(shape['fill_color'])

            if 'line' in shape and 'color' in shape['line']:
                all_colors.append(shape['line']['color'])

    # Find most common font properties
    font_families = [f['family'] for f in all_fonts if 'family' in f]
    font_sizes = [f['size'] for f in all_fonts if 'size' in f]
    font_colors = [f['color'] for f in all_fonts if 'color' in f]

    family_counter = Counter(font_families)
    size_counter = Counter(font_sizes)
    color_counter = Counter(font_colors)
    fill_counter = Counter(all_fill_colors)

    # Build style dictionary
    style = {
        'slideWidth': emu_to_inches(prs.slide_width),
        'slideHeight': emu_to_inches(prs.slide_height),
        'fonts': {},
        'colors': {},
    }

    # Most common fonts by size (assuming larger = title, smaller = body)
    if size_counter:
        sorted_sizes = sorted(size_counter.items(), key=lambda x: x[1], reverse=True)

        # Title font (largest common size)
        large_sizes = [s for s, c in sorted_sizes if s >= 24]
        if large_sizes:
            style['fonts']['title'] = {
                'size': max(large_sizes),
                'family': family_counter.most_common(1)[0][0] if family_counter else 'Arial'
            }

        # Subtitle font (medium size)
        medium_sizes = [s for s, c in sorted_sizes if 16 <= s < 24]
        if medium_sizes:
            style['fonts']['subtitle'] = {
                'size': medium_sizes[0] if medium_sizes else 18,
                'family': family_counter.most_common(1)[0][0] if family_counter else 'Arial'
            }

        # Body font (smaller size)
        small_sizes = [s for s, c in sorted_sizes if s < 16]
        if small_sizes:
            style['fonts']['body'] = {
                'size': small_sizes[0] if small_sizes else 11,
                'family': family_counter.most_common(1)[0][0] if family_counter else 'Arial'
            }

        # Source/footnote font (smallest size)
        if small_sizes:
            style['fonts']['source'] = {
                'size': min(small_sizes),
                'family': family_counter.most_common(1)[0][0] if family_counter else 'Arial'
            }

    # Most common colors
    if color_counter:
        common_colors = color_counter.most_common(5)
        color_names = ['darkNavy', 'mediumBlue', 'lightBlue', 'accentOrange', 'gray']
        for i, (color, count) in enumerate(common_colors):
            if i < len(color_names):
                style['colors'][color_names[i]] = color

    # Most common fill colors
    if fill_counter:
        common_fills = fill_counter.most_common(3)
        fill_names = ['primaryFill', 'secondaryFill', 'accentFill']
        for i, (color, count) in enumerate(common_fills):
            if i < len(fill_names):
                style['colors'][fill_names[i]] = color

    return style


def build_pattern_definition(pattern_name: str, slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Build a pattern definition from multiple slides of the same pattern.

    Args:
        pattern_name: Name of the pattern
        slides_data: List of slide data dictionaries for this pattern

    Returns:
        Pattern definition dictionary
    """
    if not slides_data:
        return {}

    # Get pattern info
    info = PATTERN_INFO.get(pattern_name, {})

    pattern = {
        'description': info.get('description', ''),
        'bestFor': info.get('bestFor', []),
        'elements': {},
        'exampleSlides': [s['slide_num'] for s in slides_data],
    }

    # Collect all shapes across all slides of this pattern
    all_shapes = []
    for slide_data in slides_data:
        all_shapes.extend(slide_data.get('shapes', []))

    if not all_shapes:
        return pattern

    # Categorize shapes by type and position
    shape_types = {}
    for shape in all_shapes:
        shape_type = shape.get('type', 'UNKNOWN')
        if shape_type not in shape_types:
            shape_types[shape_type] = []
        shape_types[shape_type].append(shape)

    # Build element definitions based on shape types
    for shape_type, shapes in shape_types.items():
        if len(shapes) == 0:
            continue

        # Average positions and sizes
        avg_x = sum(s['x'] for s in shapes) / len(shapes)
        avg_y = sum(s['y'] for s in shapes) / len(shapes)
        avg_width = sum(s['width'] for s in shapes) / len(shapes)
        avg_height = sum(s['height'] for s in shapes) / len(shapes)

        element_key = f"{shape_type.lower()}_{len(pattern['elements']) + 1}"

        element = {
            'type': shape_type,
            'x': round(avg_x, 2),
            'y': round(avg_y, 2),
            'width': round(avg_width, 2),
            'height': round(avg_height, 2),
            'count': len(shapes),
        }

        # Add common properties
        fill_colors = [s.get('fill_color') for s in shapes if 'fill_color' in s]
        if fill_colors:
            element['fill_color'] = Counter(fill_colors).most_common(1)[0][0]

        # Add font properties for text shapes
        all_fonts = []
        for shape in shapes:
            if 'fonts' in shape:
                all_fonts.extend(shape['fonts'])

        if all_fonts:
            font_families = [f['family'] for f in all_fonts if 'family' in f]
            font_sizes = [f['size'] for f in all_fonts if 'size' in f]

            if font_families or font_sizes:
                element['font'] = {}
                if font_families:
                    element['font']['family'] = Counter(font_families).most_common(1)[0][0]
                if font_sizes:
                    element['font']['size'] = round(sum(font_sizes) / len(font_sizes), 1)

        # Add table properties
        table_shapes = [s for s in shapes if 'table' in s]
        if table_shapes:
            table_data = table_shapes[0]['table']
            element['table'] = {
                'rows': table_data.get('rows', 0),
                'columns': table_data.get('columns', 0),
            }

        # Add chart properties
        chart_shapes = [s for s in shapes if 'chart' in s]
        if chart_shapes:
            chart_types = [s['chart'].get('chart_type') for s in chart_shapes]
            element['chart'] = {
                'type': Counter(chart_types).most_common(1)[0][0] if chart_types else 'UNKNOWN'
            }

        pattern['elements'][element_key] = element

    return pattern


def main():
    """Main function to orchestrate pattern extraction."""
    parser = argparse.ArgumentParser(
        description='Extract layout patterns from a PowerPoint template'
    )
    parser.add_argument(
        'template_path',
        nargs='?',
        help='Path to the template PPTX file'
    )
    parser.add_argument(
        '--output',
        '-o',
        help='Output JSON file path',
        default=None
    )

    args = parser.parse_args()

    # Determine template path
    if args.template_path:
        template_path = Path(args.template_path)
    else:
        # Default relative path
        script_dir = Path(__file__).parent
        template_path = script_dir / '..' / '..' / '251219_Escort_Phase 1 Market Selection_V3.pptx'

    # Determine output path
    if args.output:
        output_path = Path(args.output)
    else:
        script_dir = Path(__file__).parent
        output_path = script_dir / '..' / '..' / 'backend' / 'market-research' / 'template-patterns.json'

    # Check if template exists
    if not template_path.exists():
        logger.error(f"Template file not found: {template_path}")
        sys.exit(1)

    logger.info(f"Loading template from: {template_path}")

    try:
        # Load presentation
        prs = Presentation(str(template_path))
        logger.info(f"Template loaded successfully. Total slides: {len(prs.slides)}")

        # Extract data from all slides
        all_slides_data = []
        pattern_slides = {}

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            logger.info(f"Processing slide {slide_num}...")

            # Extract shapes
            shapes = extract_slide_shapes(slide)

            slide_data = {
                'slide_num': slide_num,
                'shapes': shapes,
            }

            # Classify slide
            pattern_name = classify_slide(slide_num)
            if pattern_name:
                slide_data['pattern'] = pattern_name
                if pattern_name not in pattern_slides:
                    pattern_slides[pattern_name] = []
                pattern_slides[pattern_name].append(slide_data)

            all_slides_data.append(slide_data)

        logger.info("Extracting global style...")
        global_style = extract_global_style(all_slides_data, prs)

        # Build pattern definitions
        logger.info("Building pattern definitions...")
        patterns = {}
        for pattern_name, slides_data in pattern_slides.items():
            logger.info(f"  Processing pattern: {pattern_name} ({len(slides_data)} slides)")
            pattern_def = build_pattern_definition(pattern_name, slides_data)
            if pattern_def:
                patterns[pattern_name] = pattern_def

        # Build final output
        output = {
            'style': global_style,
            'patterns': patterns,
            'metadata': {
                'template_file': str(template_path.name),
                'total_slides': len(prs.slides),
                'patterns_extracted': len(patterns),
            }
        }

        # Create output directory if needed
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Write output JSON
        logger.info(f"Writing output to: {output_path}")
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(output, f, indent=2, ensure_ascii=False)

        logger.info(f"Pattern extraction complete!")
        logger.info(f"Extracted {len(patterns)} patterns from {len(prs.slides)} slides")
        logger.info(f"Output saved to: {output_path}")

        # Print summary
        print("\n" + "="*60)
        print("PATTERN EXTRACTION SUMMARY")
        print("="*60)
        print(f"Template: {template_path.name}")
        print(f"Total slides: {len(prs.slides)}")
        print(f"Patterns extracted: {len(patterns)}")
        print("\nPatterns:")
        for pattern_name, pattern_def in patterns.items():
            example_count = len(pattern_def.get('exampleSlides', []))
            print(f"  - {pattern_name}: {example_count} example slides")
        print(f"\nOutput: {output_path}")
        print("="*60)

    except Exception as e:
        logger.error(f"Error processing template: {e}", exc_info=True)
        sys.exit(1)


if __name__ == "__main__":
    main()
